import logging
import pandas as pd
import backoff
import base64
from io import BytesIO
import json
import langid
import time
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed
import os
import re
from dotenv import load_dotenv
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document as LCDocument
from docling.document_converter import DocumentConverter
from docling.datamodel.base_models import InputFormat
from docling_core.types.doc import TextItem, ProvenanceItem, TableItem, PictureItem
import openai
from docling.document_converter import (
    DocumentConverter,
    PdfFormatOption,
    PowerpointFormatOption,
    FormatOption,
)
from docling.pipeline.simple_pipeline import SimplePipeline
from docling.pipeline.standard_pdf_pipeline import StandardPdfPipeline
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
from docling.backend.mspowerpoint_backend import MsPowerpointDocumentBackend
from docling.backend.docling_parse_backend import DoclingParseDocumentBackend
from docling_core.types.doc import RefItem, TextItem, BoundingBox
logging.basicConfig(
    filename='app.log',        # File to write logs to
    level=logging.DEBUG,       # Set the minimum log level
    format='%(asctime)s - %(levelname)s - %(message)s'  # Log format
)

IMAGE_RESOLUTION_SCALE = 1.0


class DoclingFileLoader(BaseLoader):
    def __init__(self, client, file_path: str | list[str]) -> None:
        load_dotenv()
        self.bedrock_client = client
        self._file_paths = file_path if isinstance(file_path, list) else [file_path]
        self.pipeline_options = PdfPipelineOptions()
        self.pipeline_options.do_ocr = False
        self.pipeline_options.ocr_options.use_gpu = False
        self.pipeline_options.do_table_structure = True
        self.pipeline_options.images_scale = IMAGE_RESOLUTION_SCALE
        self.pipeline_options.generate_picture_images = True
        self.doc_converter = DocumentConverter(
            allowed_formats=[InputFormat.PPTX, InputFormat.PDF, InputFormat.IMAGE],
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=self.pipeline_options,
                    backend=PyPdfiumDocumentBackend,
                ),
                InputFormat.PPTX: PowerpointFormatOption(
                    pipeline_cls=SimplePipeline,
                    backend=MsPowerpointDocumentBackend,
                ),
                InputFormat.IMAGE: FormatOption(
                    pipeline_cls=StandardPdfPipeline,
                    backend=DoclingParseDocumentBackend,
                ),
            },
        )

    def detect_language(self, doc):
        if doc.texts:
            text_samples = [item.text for item in doc.texts[:25]]
            combined_text = " ".join(text_samples)
            
            try:
                # Use langid for language detection
                lang, _ = langid.classify(combined_text)
                return lang
            except Exception as e:
                print(f"Error detecting language: {e}")
                logging.warning(f"Error detecting language: {e}")
                return 'en'  # Default to English if detection fails
        return 'en'  # Default to English if no text is available

    def process_image(self, image_data, doc_language):
        if isinstance(image_data, bytes):
            image_data = base64.b64encode(image_data).decode('utf-8')

        try:
            # Prepare the body for the model invocation
            body = json.dumps(
                {
                    "anthropic_version": "bedrock-2023-05-31",
                    "max_tokens": 500,
                    "messages": [
                        {
                            "role": "user",
                            "content": json.dumps(
                                {
                                    "type": "composite",
                                    "elements": [
                                        {
                                            "type": "image",
                                            "source": {
                                                "type": "base64",
                                                "format": "jpeg",
                                                "media_type": "image/jpeg",
                                                "data": image_data[:1000],  # Ensure 'data' contains the Base64 string
                                            },
                                        },
                                        {
                                            "type": "text",
                                            "text": f"Provide a clear and detailed description of the contents of the image. The language of the provided description should be {doc_language}.",
                                        },
                                    ],
                                }
                            ),
                        }
                    ],
                }
            )
            @backoff.on_exception(backoff.expo, openai.RateLimitError, max_tries=3)
            def get_response(modelId, body, max_tokens=500):
                return self.bedrock_client.invoke_model(modelId=modelId, body=body)
 
            
            try:
                # Call the model
                response = get_response("anthropic.claude-3-sonnet-20240229-v1:0", body)
            except openai.RateLimitError as e:
                print("Error: Rate limit exceeded. Please try again later.")
                logging.error("Error: Rate limit exceeded. Please try again later.")
                time.sleep(60) 
            

            response_body = json.loads(response.get("body").read())
        

            return response_body['content'][0]['text']

        except Exception as e:
            print(f"Error in generate_image_description: {e}")
            logging.error(f"Error in generate_image_description: {e}")


    def process_pptx_page_images(self, dlc_doc, pptx_path) -> None:
        """Process the images of figures and tables from a PPTX document and collect image data concurrently."""
        figure_images = []  # To store images of figures and their slide numbers

        # Load the presentation
        prs = Presentation(pptx_path)

        # Loop through slides and extract images
        for slide_ind, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Type 13 corresponds to Picture shapes
                    # Extract image data (use .image for images)
                    img_stream = shape.image
                    img_format = img_stream.ext  # Format of the image (e.g., 'jpeg', 'png')
                    img_data = img_stream.blob  # Raw image data
                    img_base64 = base64.b64encode(img_data).decode('utf-8')  # Base64 encode the image

                    # Store the base64 image data along with the slide number (page number equivalent)
                    figure_images.append((slide_ind + 1, img_base64))  # slide_ind + 1 to make it 1-indexed

        # Initialize ThreadPoolExecutor for concurrent processing
        doc_language = self.detect_language(dlc_doc)  # Assuming you have a method to detect language
        descriptions = {}
        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = {}

            # Submit the image processing tasks concurrently
            for page_no, img_base64 in figure_images:
                futures[executor.submit(self.process_image, img_base64, doc_language)] = page_no

            # Process the results as tasks complete
            for future in as_completed(futures):
                try:
                    description = future.result()
                    page_no = futures[future]

                    if description:
                        descriptions[page_no] = description
                except Exception as e:
                    print(f"Error processing image: {e}")
                    logging.error(f"Error processing image: {e}")

        # Interleave descriptions into the document at the correct figure locations
        self.interleave_descriptions(dlc_doc, descriptions)  

    # try replacing it with a different PDF image extraction library
    def process_pdf_page_images(self, dl_doc) -> None:
        """Process the images of figures and tables from the document and collect image data concurrently."""
        figure_images = []  # To store images of tables and figures

        # Collect all figure images and their page numbers
        for element, _level in dl_doc.iterate_items():
            if isinstance(element, PictureItem):
                # Extract page_no from the ProvenanceItem
                page_no = element.prov[0].page_no if element.prov else None

                # Extract image data
                pil_image = element.image.pil_image  # PIL.Image object
                img_byte_arr = BytesIO()
                pil_image.save(img_byte_arr, format="JPEG")  # Convert to JPEG format
                img_byte_arr = img_byte_arr.getvalue()  # Get byte data from byte stream
                img_base64 = base64.b64encode(img_byte_arr).decode("utf-8")  # Convert image byte data to base64
                
                # Store the image base64 data along with the page number
                figure_images.append((page_no, img_base64))

        # Initialize ThreadPoolExecutor for concurrent processing
        doc_language = self.detect_language(dl_doc)
        descriptions = {}
        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = {}

            # Submit the image processing tasks concurrently
            for page_no, img_base64 in figure_images:
                futures[executor.submit(self.process_image, img_base64, doc_language)] = page_no

            # Process the results as tasks complete
            for future in as_completed(futures):
                try:
                    description = future.result()
                    page_no = futures[future]

                    if description:
                        descriptions[page_no] = description
                except Exception as e:
                    print(f"Error processing image: {e}")
                    logging.error(f"Error processing image: {e}")

        # Interleave descriptions into the document at the correct figure locations
        self.interleave_descriptions(dl_doc, descriptions)


    def interleave_descriptions(self, dl_doc, descriptions):
        """Interleave figure descriptions into the document at the correct locations"""
        default_bbox = BoundingBox(l=0, t=0, r=100, b=100)

        for page_no, description in descriptions.items():
            # Find the correct place to insert the description based on the page number
            page_text_items = [item for item in dl_doc.texts if item.prov[0].page_no == page_no]

            if page_text_items:
                last_text_item = page_text_items[-1]  # The last text item for this page
                text_id = len(dl_doc.texts)  # Get a unique text ID
                self_ref = f"#/texts/{text_id}"

                # Create a new TextItem for the description
                new_text_item = TextItem(
                    self_ref=self_ref,
                    parent=RefItem(cref="#/body"),
                    label="paragraph",
                    prov=[ProvenanceItem(page_no=page_no, bbox=default_bbox, coord_origin=None, charspan=(0, len(description)))],
                    orig=description,
                    text=description,
                )

                # Insert this new TextItem right after the last text item for this page
                last_index = dl_doc.texts.index(last_text_item)
                dl_doc.texts.insert(last_index + 1, new_text_item)
            else:
                # If no text items exist for this page, just append the description
                text_id = len(dl_doc.texts)
                self_ref = f"#/texts/{text_id}"

                # Create the new TextItem with the description
                new_text_item = TextItem(
                    self_ref=self_ref,
                    parent=RefItem(cref="#/body"),  # Parent can be '#/body' or wherever needed
                    label="paragraph",
                    prov=[ProvenanceItem(page_no=page_no, bbox=default_bbox, coord_origin=None, charspan=(0, len(description)))],
                    orig=description,
                    text=description,
                )

                dl_doc.texts.append(new_text_item)
              


    def lazy_load(self) -> list[LCDocument]:
        documents = []  # Initialize an empty list to store documents
        output_file_path = "docling_document_structure.txt"  # Define the output file path

        with open(output_file_path, 'w') as output_file:  # Open the file for writing
            for source in self._file_paths:
                try:
                    print(f"Processing {source}")
                    logging.info(f"Processing {source}")
                    dl_doc = self.doc_converter.convert(source).document
                    # Check if the document is a PDF before processing images
                    if source.lower().endswith('.pdf'):
                        print("Processing PDF page images...")
                        logging.info("Processing PDF page images...")
                        self.process_pdf_page_images(dl_doc)
                    elif source.lower().endswith('.pptx'):
                        print("Processing PPTX page images...")
                        logging.info("Processing PPTX page images...")
                        self.process_pptx_page_images(dl_doc, source)
                    
                    # Write the structure of the DoclingDocument to the output file
                    output_file.write(f"Docling Document structure for {source}:\n")
                    output_file.write(f"{dl_doc.__dict__}\n\n")  # Write the dictionary structure
                    print(f"Docling Document structure for {source} written to file.")
                    logging.info(f"Docling Document structure for {source} written to file.")
                    file_name = os.path.basename(source)
                 
                    with open('source_data.txt', "r") as file:
                    
                        for line in file:
                            if f"filename: {file_name}" in line:
                             
                                # Split the line and extract the source data
                                parts = line.split(", source: ")
                                source_url = parts[1].strip()  # Return the source, stripping any extra whitespace
                                match = re.match(r"(.*?/sites/[^/]+)", source_url)
                                if match:
                                    extracted_url = match.group(1)  # Extract the matched portion
                                    # Append the filename
                                    source_url = f"{extracted_url}/{file_name}"
                                    dl_doc.print_element_tree()
                                    for item, level in dl_doc.iterate_items():
                                        if isinstance(item, TextItem):
                                            for prov in item.prov: 
                                                text = item.text
                                                page_no = prov.page_no
                                                documents.append(LCDocument(page_content=text, metadata={"source": source_url, "filename": file_name, "page_number" : page_no}))
                                        elif isinstance(item, TableItem):
                                            table_df: pd.DataFrame = item.export_to_dataframe()
                                            text = table_df.to_markdown()
                                            for prov in item.prov:
                                                page_no = prov.page_no
                                                documents.append(LCDocument(page_content=text, metadata={"source": source_url, "filename": file_name, "page_number" : page_no}))
                                    '''if dl_doc.tables:
                                        for table_ix, table in enumerate(dl_doc.tables):
                                            table_df: pd.DataFrame = table.export_to_dataframe()
                                            print(f"## Table {table_ix}")
                                            text = table_df.to_markdown()
                                            documents.append(LCDocument(page_content=text, metadata={"source": source_url, "filename": file_name}))
                                    else:
                                        # Export text if there are no images
                                        text = dl_doc.export_to_markdown()
                                        documents.append(LCDocument(page_content=text, metadata={"source": source_url, "filename": file_name}))'''  # Append the document to the list
                                    
                except PermissionError as e:
                    print(f"PermissionError: {e}")
                    logging.error(f"PermissionError: {e}")
                    print("Ensure that the file is not open in another program and you have the necessary permissions.")
                    logging.warning("Ensure that the file is not open in another program and you have the necessary permissions.")
                    continue  # Skip this file and continue with others
                except FileNotFoundError as e:
                    print(f"FileNotFoundError: {e}")
                    logging.error(f"FileNotFoundError: {e}")
                    print("Ensure that the file exists at the specified path.")
                    logging.warning("Ensure that the file exists at the specified path.")
                    continue  # Skip this file and continue with others
                except Exception as e:
                    print(f"Unexpected error processing {source}: {e}")
                    logging.error(f"Unexpected error processing {source}: {e}")
                    continue  # Skip this file and continue with others

        return documents  # Return the list of documents


    def load(self):
        try:
            return list(self.lazy_load())  # Will consume the generator
        except StopIteration:
            return []  # Handle StopIteration gracefully, return empty list
        except Exception as e:
            print(f"An error occurred: {e}")
            logging.error(f"An error occurred: {e}")
            return []
